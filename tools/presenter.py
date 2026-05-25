"""tools/presenter.py — Intelligent AI-driven presentation assistant.

JARVIS finds the file intelligently, opens it in slideshow mode, and narrates
each slide with LLM-powered explanation — not just reading text, but understanding
and presenting the content like a real presenter.
"""

import os
import subprocess
import threading
import time
from typing import Optional

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

_session:     dict = {}
_speak_fn     = None
_tts_ref      = None   # TTSEngine — used to wait for speech to finish
_pause_hook   = None   # callable() — called when presentation starts (pause bg services)
_resume_hook  = None   # callable() — called when presentation ends  (resume bg services)

_SEARCH_ROOTS = [
    os.path.expanduser("~/Desktop"),
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"),
    os.path.expanduser("~/OneDrive"),
    os.path.expanduser("~/OneDrive - Personal"),
    os.path.expanduser("~"),
]

_PPT_EXE_PATHS = [
    r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\root\Office16\POWERPNT.EXE",
    r"C:\Program Files\Microsoft Office\Office16\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\Office16\POWERPNT.EXE",
    r"C:\Program Files\Microsoft Office\root\Office15\POWERPNT.EXE",
]


def _set_speak(fn) -> None:
    global _speak_fn
    _speak_fn = fn


def _set_tts(tts) -> None:
    global _tts_ref
    _tts_ref = tts


def _set_busy_hooks(pause_fn, resume_fn) -> None:
    """Register callbacks invoked when presentation starts/ends (pause background services)."""
    global _pause_hook, _resume_hook
    _pause_hook  = pause_fn
    _resume_hook = resume_fn


# ------------------------------------------------------------------ #
#  Auto-advance engine                                                #
# ------------------------------------------------------------------ #

def _wait_tts_done(min_pause: float = 0.5) -> None:
    """Block until TTS is fully silent, then add a breath pause."""
    time.sleep(min_pause)
    if _tts_ref:
        while _tts_ref.is_speaking():
            time.sleep(0.1)
    time.sleep(0.5)


def _speak_narration(text: str, slide_idx: int) -> None:
    """Speak narration while printing a heartbeat every 1.5 s.

    The heartbeat keeps automation's speech_idle counter below the 2.2 s
    threshold so observe_turn doesn't exit mid-narration.
    """
    if not _speak_fn:
        return
    _stop_ticker = threading.Event()

    def _ticker():
        total = _session.get("total", "?")
        while not _stop_ticker.wait(timeout=1.5):
            print(f"[PresentAuto] Narrating slide {slide_idx + 1} of {total}...", flush=True)

    t = threading.Thread(target=_ticker, daemon=True, name="PresentTicker")
    t.start()
    try:
        _speak_fn(text)
    finally:
        _stop_ticker.set()


def _focus_ppt_window() -> None:
    """Bring the PowerPoint slideshow to the foreground so key presses land on it."""
    try:
        import pygetwindow as gw
        wins = [w for w in gw.getAllWindows()
                if w.title and ("PowerPoint" in w.title
                                or ".pptx" in w.title.lower()
                                or ".ppt"  in w.title.lower())]
        if wins:
            wins[0].activate()
            time.sleep(0.2)
    except Exception:
        pass


def _present_sequence(narration_0: str = "", skip_initial: bool = False) -> None:
    """Background thread: [wait for intro TTS →] narrate slide 0 → auto-advance slides 1+.

    skip_initial=True: skip the slide-0 narration and go straight to auto-advance
    from the current slide index (used by resume_presentation).
    """
    try:
        if not skip_initial:
            # Wait for brain to finish speaking the short "Opening…" intro message
            time.sleep(1.5)
            if _tts_ref:
                while _tts_ref.is_speaking():
                    time.sleep(0.1)
            time.sleep(0.4)

            # Narrate slide 0 directly — PPT is on slide 1, display matches narration
            if not _session.get("auto_stopped") and _speak_fn:
                _speak_narration(narration_0, 0)
            # Signal automation that slide 1 is done so observe_turn can exit cleanly.
            if not _session.get("auto_stopped"):
                print("[PresentAuto] Slide 1 narration complete.", flush=True)

        # Auto-advance slides 1, 2, 3, …  (or resume from current when skip_initial)
        while not _session.get("auto_stopped"):
            slides = _session.get("slides")
            if not slides:
                print("[PresentAuto] Session cleared — stopping.")
                break

            idx   = _session.get("current", 0)
            total = _session.get("total", 0)

            if idx >= total - 1:
                if not _session.get("auto_stopped") and _speak_fn:
                    _speak_fn("That's the last slide. Presentation complete.")
                _session["auto_stopped"] = True
                break

            if _session.get("auto_stopped"):
                break

            _session["current"] = idx + 1
            _focus_ppt_window()
            _press_key("right")
            time.sleep(0.4)

            if _session.get("auto_stopped"):
                break

            try:
                narration = _narrate(_session["current"])
            except Exception as e:
                print(f"[PresentAuto] Narrate error slide {_session['current']}: {e}")
                narration = _narrate_slide_raw(
                    _session["slides"][_session["current"]],
                    _session["current"], total
                )

            # Double gate: also check slides — _session.clear() removes auto_stopped
            # so we can't rely on it alone after end_presentation() clears the dict.
            if _speak_fn and _session.get("slides") and not _session.get("auto_stopped"):
                _speak_narration(narration, _session["current"])

            if _session.get("auto_stopped") or not _session.get("slides"):
                break

            _wait_tts_done()

    except Exception as e:
        print(f"[PresentAuto] Fatal error — presentation stopped: {e}")
        import traceback
        traceback.print_exc()


# ------------------------------------------------------------------ #
#  Multi-provider LLM for intelligent narration                       #
#  Uses requests directly — no SDK packages required.                 #
# ------------------------------------------------------------------ #

def _quick_llm(prompt: str, system: str = "") -> str:
    """Call an LLM for slide narration. Tries all configured providers in order."""
    import requests as _req

    msgs: list = []
    if system:
        msgs.append({"role": "system", "content": system})
    msgs.append({"role": "user", "content": prompt})

    _PROVIDERS = [
        # (env_key, url, model, extra_headers)
        (
            "GROQ_API_KEY",
            "https://api.groq.com/openai/v1/chat/completions",
            "llama-3.3-70b-versatile",
            {},
        ),
        (
            "GEMINI_API_KEY",
            "https://generativelanguage.googleapis.com/v1beta/openai/chat/completions",
            "gemini-1.5-flash",
            {},
        ),
        (
            "NVIDIA_API_KEY",
            "https://integrate.api.nvidia.com/v1/chat/completions",
            "meta/llama-3.1-70b-instruct",
            {},
        ),
        (
            "OPENROUTER_API_KEY",
            "https://openrouter.ai/api/v1/chat/completions",
            "mistralai/mistral-7b-instruct:free",
            {"HTTP-Referer": "https://github.com/jarvis-assistant"},
        ),
    ]

    for env_key, url, model, extra_hdrs in _PROVIDERS:
        api_key = os.environ.get(env_key, "")
        if not api_key:
            continue
        try:
            hdrs = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type":  "application/json",
                **extra_hdrs,
            }
            resp = _req.post(
                url,
                headers=hdrs,
                json={"model": model, "messages": msgs,
                      "max_tokens": 300, "temperature": 0.7, "stream": False},
                timeout=12,
            )
            resp.raise_for_status()
            text = resp.json()["choices"][0]["message"]["content"].strip()
            if text:
                return text
        except Exception as e:
            print(f"[Presenter] {env_key} failed: {e}")

    return ""


# ------------------------------------------------------------------ #
#  Smart file discovery                                               #
# ------------------------------------------------------------------ #

def _all_ppt_files_in(root: str, max_depth: int = 4) -> list[str]:
    """Walk a directory up to max_depth and collect all pptx/ppt/pdf files."""
    results = []
    if not os.path.isdir(root):
        return results
    root_depth = root.rstrip(os.sep).count(os.sep)
    try:
        for dirpath, dirnames, filenames in os.walk(root):
            if dirpath.count(os.sep) - root_depth >= max_depth:
                dirnames.clear()
                continue
            # Skip hidden/system dirs
            dirnames[:] = [d for d in dirnames if not d.startswith('.') and d not in ('$RECYCLE.BIN', 'Windows', 'System32')]
            for f in filenames:
                if f.lower().endswith(('.pptx', '.ppt', '.pdf')):
                    results.append(os.path.join(dirpath, f))
    except PermissionError:
        pass
    return results


def _find_ppt_smart(query: str) -> Optional[str]:
    """
    Resolve a query to a real file path using multiple strategies:
    1. Exact absolute path
    2. Path relative to common roots
    3. Filename-only search across common roots (with keyword matching)
    4. Fuzzy keyword match across all pptx/ppt files in common roots
    """
    query = os.path.expandvars(query.strip().strip('"').strip("'").strip())

    # Strategy 1: absolute path or close to it
    if os.path.exists(query):
        return os.path.abspath(query)
    # Try adding common extensions
    for ext in ('.pptx', '.ppt', '.pdf'):
        if not query.lower().endswith(ext) and os.path.exists(query + ext):
            return os.path.abspath(query + ext)

    # Strategy 2: relative to common roots
    for root in _SEARCH_ROOTS:
        if not os.path.isdir(root):
            continue
        candidate = os.path.join(root, query)
        if os.path.exists(candidate):
            return candidate
        for ext in ('.pptx', '.ppt', '.pdf'):
            if os.path.exists(candidate + ext):
                return candidate + ext
        # Try just the basename
        basename_candidate = os.path.join(root, os.path.basename(query))
        if os.path.exists(basename_candidate):
            return basename_candidate

    # Strategy 3 & 4: keyword search
    # Extract keywords from query (strip extension words, short words)
    stopwords = {'the', 'a', 'an', 'my', 'this', 'that', 'file', 'ppt',
                 'pptx', 'pdf', 'presentation', 'present', 'show', 'open'}
    raw_keywords = query.replace('.pptx', '').replace('.ppt', '').replace('.pdf', '').split()
    keywords = [k.lower() for k in raw_keywords if len(k) > 2 and k.lower() not in stopwords]

    all_files: list[str] = []
    for root in _SEARCH_ROOTS:
        all_files.extend(_all_ppt_files_in(root))

    if not all_files:
        return None

    # Score each file: how many keywords appear in its name?
    def score(filepath: str) -> int:
        name = os.path.basename(filepath).lower()
        return sum(1 for k in keywords if k in name)

    if keywords:
        scored = [(score(f), f) for f in all_files if score(f) > 0]
        if scored:
            scored.sort(key=lambda x: (-x[0], -os.path.getmtime(x[1])))
            return scored[0][1]

    # Last resort: return the most recently modified pptx/ppt if no keywords matched
    ppt_only = [f for f in all_files if f.lower().endswith(('.pptx', '.ppt'))]
    if ppt_only:
        ppt_only.sort(key=os.path.getmtime, reverse=True)
        return ppt_only[0]

    return None


# ------------------------------------------------------------------ #
#  Text extraction from slides                                        #
# ------------------------------------------------------------------ #

def _extract_slide_text(slide) -> dict:
    """Extract ALL text from a slide — title (by font size / placeholder), bullets, notes."""
    items = []  # {y, font_pt, lines, is_title_ph}

    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            continue
        if not hasattr(shape, "text_frame"):
            continue

        paras = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
        if not paras:
            continue

        # Check for standard title placeholder (idx 0 or 13)
        is_title_ph = False
        try:
            ph = shape.placeholder_format
            if ph and ph.idx in (0, 13):
                is_title_ph = True
        except Exception:
            pass

        # Find max font size in the shape
        max_pt = 0.0
        for para in shape.text_frame.paragraphs:
            try:
                if para.font.size:
                    max_pt = max(max_pt, para.font.size.pt)
            except Exception:
                pass
            for run in para.runs:
                try:
                    if run.font.size:
                        max_pt = max(max_pt, run.font.size.pt)
                except Exception:
                    pass

        y = getattr(shape, "top", 999_999_999) or 999_999_999
        items.append({"y": y, "font_pt": max_pt, "lines": paras, "is_title_ph": is_title_ph})

    if not items:
        notes = _extract_notes(slide)
        return {"title": "", "bullets": [], "notes": notes}

    # Title = explicit title placeholder > largest font > topmost shape
    title_item = next((it for it in items if it["is_title_ph"]), None)
    if not title_item:
        sorted_by_size = sorted(items, key=lambda x: (-x["font_pt"], x["y"]))
        title_item = sorted_by_size[0]

    title = " ".join(title_item["lines"])
    bullets = [line for it in items if it is not title_item for line in it["lines"]]

    return {"title": title, "bullets": bullets, "notes": _extract_notes(slide)}


def _extract_notes(slide) -> str:
    try:
        ns = slide.notes_slide
        if ns:
            tf = ns.notes_text_frame
            if tf:
                return " ".join(p.text.strip() for p in tf.paragraphs if p.text.strip())
    except Exception:
        pass
    return ""


# ------------------------------------------------------------------ #
#  LLM-powered narration                                              #
# ------------------------------------------------------------------ #

def _llm_narrate(slide_data: dict, slide_num: int, total: int,
                 presentation_context: str, prev_topic: str = "") -> str:
    """Intelligently narrate a slide using whatever LLM is available."""
    parts = []
    if slide_data["title"]:
        parts.append(f"Title: {slide_data['title']}")
    if slide_data["bullets"]:
        parts.append("Content:\n" + "\n".join(f"- {b}" for b in slide_data["bullets"]))
    if slide_data["notes"]:
        parts.append(f"Speaker notes: {slide_data['notes']}")

    if not parts:
        return f"Slide {slide_num} of {total} appears to have no text — it may be an image or diagram."

    prev_ctx = f"\nPrevious slide covered: {prev_topic}." if prev_topic else ""
    system = (
        "You are an intelligent presentation assistant. "
        "Given slide content, explain it naturally and insightfully — "
        "not just reading the text, but understanding the meaning, adding context, "
        "and presenting like a real speaker. 2-4 sentences. "
        "Do not start with 'This slide'."
    )
    prompt = (
        f"Presentation: {presentation_context}\n"
        f"Slide {slide_num} of {total}.{prev_ctx}\n\n"
        + "\n".join(parts)
        + "\n\nPresent this slide naturally."
    )
    return _quick_llm(prompt, system)


def _narrate_slide_raw(data: dict, idx: int, total: int) -> str:
    """Fallback: plain text narration if LLM unavailable."""
    parts: list[str] = []
    if data["title"]:
        parts.append(f"Slide {idx + 1} of {total}: {data['title']}.")
    else:
        parts.append(f"Slide {idx + 1} of {total}.")
    for bullet in data["bullets"]:
        b = bullet.strip()
        if b and b[-1] not in ".!?:":
            b += "."
        if b:
            parts.append(b)
    if data.get("notes"):
        parts.append(f"Speaker notes: {data['notes']}.")
    return " ".join(parts) if parts else f"Slide {idx + 1} has no text content."


def _narrate(idx: int) -> str:
    """Narrate slide at index idx — LLM first, fallback to plain text."""
    data = _session["slides"][idx]
    total = _session["total"]
    context = _session.get("context", "")
    prev = _session.get("prev_topic", "")

    narration = _llm_narrate(data, idx + 1, total, context, prev)
    if not narration:
        narration = _narrate_slide_raw(data, idx, total)

    # Store topic summary for next slide's context
    _session["prev_topic"] = data["title"] or (data["bullets"][0] if data["bullets"] else "")
    return narration


# ------------------------------------------------------------------ #
#  PowerPoint launcher                                                #
# ------------------------------------------------------------------ #

def _find_powerpnt() -> Optional[str]:
    for p in _PPT_EXE_PATHS:
        if os.path.exists(p):
            return p
    return None


def _open_slideshow(path: str) -> Optional[str]:
    """Open in PowerPoint slideshow mode (/S flag). Returns error string or None."""
    ppt_exe = _find_powerpnt()
    if ppt_exe:
        try:
            subprocess.Popen([ppt_exe, "/S", path])
            return None
        except Exception as e:
            pass  # fall through to startfile
    try:
        os.startfile(path)
        return None
    except Exception as e:
        return str(e)


def _press_key(key: str) -> None:
    try:
        import pyautogui
        pyautogui.press(key)
    except Exception:
        pass


# ------------------------------------------------------------------ #
#  Public tools                                                        #
# ------------------------------------------------------------------ #

def present_file(path: str) -> str:
    """
    Find and present a PowerPoint or PDF. Accepts:
    - Full path: C:\\Users\\...\\file.pptx
    - Filename only: TCS_Townhall.pptx
    - Keywords: TCS townhall  (searches Desktop, Documents, Downloads, OneDrive)
    Opens in slideshow mode and narrates intelligently using AI.
    """
    if not path or not path.strip():
        return "Please tell me the name or path of the presentation to open."

    resolved = _find_ppt_smart(path)
    if not resolved:
        return (
            f"I couldn't find a presentation matching '{path}'. "
            "Check that the file exists in Desktop, Documents, or Downloads, "
            "or give me the full path."
        )

    ext = os.path.splitext(resolved)[1].lower()
    filename = os.path.basename(resolved)

    if ext in (".pptx", ".ppt"):
        if not _HAS_PPTX:
            return "python-pptx is not installed. Run: pip install python-pptx"
        try:
            prs = _Presentation(resolved)
            slides_data = [_extract_slide_text(s) for s in prs.slides]
            if not slides_data:
                return f"The file '{filename}' appears to have no slides."

            # Build presentation context from filename + first slide title
            first_title = slides_data[0]["title"]
            base_name = os.path.splitext(filename)[0].replace("_", " ").replace("-", " ")
            context = first_title if first_title else base_name

            _session.clear()
            _session.update({
                "slides":       slides_data,
                "current":      0,
                "total":        len(slides_data),
                "path":         resolved,
                "type":         "pptx",
                "context":      context,
                "prev_topic":   "",
                "auto_stopped": False,
            })

            if _pause_hook:
                try:
                    _pause_hook()
                except Exception:
                    pass

            open_err = _open_slideshow(resolved)
            time.sleep(3.0)   # let PowerPoint finish opening
            _focus_ppt_window()
            _press_key("home")  # force slide 1 regardless of last-saved position
            time.sleep(0.3)

            narration_0 = _narrate(0)

            # Background thread: waits for brain's intro TTS, then narrates slide 1,
            # then auto-advances. This keeps display and narration in sync.
            threading.Thread(
                target=_present_sequence,
                kwargs={"narration_0": narration_0},
                daemon=True, name="PresentAuto",
            ).start()

            if open_err:
                return f"Found '{filename}' but couldn't open PowerPoint: {open_err}. Starting narration."
            return f"Opening '{filename}' — {len(slides_data)} slides."
        except Exception as e:
            return f"Error loading '{filename}': {e}"

    elif ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(resolved)
            pages_text = [doc[i].get_text() for i in range(len(doc))]
            doc.close()
            _session.clear()
            _session.update({
                "slides": [{"title": f"Page {i+1}", "bullets": [t.strip()[:600]], "notes": ""} for i, t in enumerate(pages_text)],
                "current": 0,
                "total": len(pages_text),
                "path": resolved,
                "type": "pdf",
                "context": os.path.splitext(filename)[0].replace("_", " "),
                "prev_topic": "",
            })
            _session["auto_stopped"] = False
            open_err = _open_slideshow(resolved)
            time.sleep(1.5)
            narration_0 = _narrate(0)
            threading.Thread(
                target=_present_sequence,
                kwargs={"narration_0": narration_0},
                daemon=True, name="PresentAuto",
            ).start()
            if open_err:
                return f"PDF loaded ({len(pages_text)} pages) but couldn't open viewer: {open_err}. Starting narration."
            return f"PDF open — {len(pages_text)} pages."
        except ImportError:
            open_err = _open_slideshow(resolved)
            if open_err:
                return f"Cannot open PDF: {open_err}"
            return f"PDF '{filename}' opened. Install PyMuPDF for narration: pip install pymupdf"

    else:
        open_err = _open_slideshow(resolved)
        if open_err:
            return f"Cannot open '{filename}': {open_err}"
        return f"Opened '{filename}'."


def pause_presentation() -> str:
    """Pause auto-advance. Narration stops immediately."""
    if not _session.get("slides"):
        return "No presentation is active."
    _session["auto_stopped"] = True
    if _tts_ref:
        try:
            _tts_ref.stop_immediately()
        except Exception:
            pass
    return "Paused. Say 'resume presentation' to continue or 'next slide' to step manually."


def resume_presentation() -> str:
    """Resume auto-advance from the current slide."""
    if not _session.get("slides"):
        return "No presentation is active."
    _session["auto_stopped"] = False
    threading.Thread(
        target=_present_sequence,
        kwargs={"skip_initial": True},
        daemon=True, name="PresentAuto",
    ).start()
    return "Resuming presentation."


def next_slide() -> str:
    """Skip to the next slide (also pauses auto-advance)."""
    if not _session.get("slides"):
        return "No presentation is active. Tell me which file to open."
    _session["auto_stopped"] = True   # pause auto so we don't double-narrate
    if _session["current"] >= _session["total"] - 1:
        return "Already on the last slide. Say 'stop presentation' to exit."
    _session["current"] += 1
    _press_key("right")
    time.sleep(0.3)
    return _narrate(_session["current"])


def prev_slide() -> str:
    """Go back to the previous slide (also pauses auto-advance)."""
    if not _session.get("slides"):
        return "No presentation is active."
    _session["auto_stopped"] = True
    if _session["current"] <= 0:
        return "Already on the first slide."
    _session["current"] -= 1
    _press_key("left")
    time.sleep(0.3)
    return _narrate(_session["current"])


def goto_slide(number: int) -> str:
    """Jump to a specific slide number and present it."""
    if not _session.get("slides"):
        return "No presentation is active."
    idx = max(0, min(int(number) - 1, _session["total"] - 1))
    _session["current"] = idx
    return _narrate(idx)


def read_current_slide() -> str:
    """Re-present the current slide with full AI narration."""
    if not _session.get("slides"):
        return "No presentation is active."
    return _narrate(_session["current"])


def presentation_overview() -> str:
    """Summarize the entire presentation — what it's about and what each slide covers."""
    if not _session.get("slides"):
        return "No presentation is loaded."
    titles = []
    for i, s in enumerate(_session["slides"]):
        t = s["title"] or (s["bullets"][0][:60] if s["bullets"] else f"Slide {i+1}")
        titles.append(f"{i+1}. {t}")

    outline = "\n".join(titles)
    summary = _quick_llm(
        f"Presentation: {_session.get('context', '')}\n\nSlide outline:\n{outline}\n\nGive a 3-sentence overview of what this presentation covers.",
        "You are summarizing a presentation for a listener. Be concise.",
    )
    if summary:
        return f"{summary} It has {_session['total']} slides."

    return f"Presentation: {_session.get('context', 'Unknown')} — {_session['total']} slides:\n" + "\n".join(titles)


def end_presentation() -> str:
    """Stop and close the presentation."""
    _session["auto_stopped"] = True   # signal background thread to exit
    if _tts_ref:
        try:
            _tts_ref.stop_immediately()
        except Exception:
            pass
    if _resume_hook:
        try:
            _resume_hook()
        except Exception:
            pass
    name = os.path.basename(_session.get("path", "")) if _session else ""
    _session.clear()
    # Re-set after clear: _session.clear() removes auto_stopped, which makes
    # `not _session.get("auto_stopped")` evaluate True — the background thread
    # would then speak the next slide even though we just stopped. Keeping this
    # flag alive prevents any pending narration from firing.
    _session["auto_stopped"] = True
    _press_key("escape")
    return f"Presentation{(' — ' + name) if name else ''} ended."
